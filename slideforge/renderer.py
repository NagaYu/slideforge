"""Renders parsed slides into a .pptx file using python-pptx.

The renderer is the only module that imports python-pptx.  It consumes:
* the slide model from :mod:`slideforge.parser`
* layout decisions / geometry from :mod:`slideforge.layout_engine`
* colors & fonts from :mod:`slideforge.themes`
* safe font sizing from :mod:`slideforge.autofit`
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from . import autofit
from . import layout_engine as le
from .layout_engine import Rect
from .parser import Block, Slide

# base font sizes (pt) -- autofit may scale these down, never up
SIZE = {
    "hero_title": 44,
    "hero_sub": 18,
    "kicker": 14,
    "title": 32,
    "col_header": 20,
    "card_title": 17,
    "body": 16,
    "card_body": 13,
    "step_label": 14,
    "footer": 10,
}


def _rgb(t) -> RGBColor:
    return RGBColor(*t)


class Renderer:
    def __init__(self, theme: dict):
        self.theme = theme
        self.c = theme["colors"]
        self.f = theme["fonts"]
        self.rules = theme["rules"]
        self.prs = Presentation()
        self.prs.slide_width = Inches(le.SLIDE_W)
        self.prs.slide_height = Inches(le.SLIDE_H)
        self.blank = self.prs.slide_layouts[6]

    # -- low level helpers -------------------------------------------------

    def _new_slide(self, dark: bool = False):
        slide = self.prs.slides.add_slide(self.blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(self.c["bg_dark"] if dark else self.c["bg"])
        return slide

    def _textbox(self, slide, rect: Rect, anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(
            Inches(rect.left), Inches(rect.top),
            Inches(rect.width), Inches(rect.height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Inches(0.02)
        tf.margin_top = tf.margin_bottom = Inches(0.02)
        return tf

    def _para(self, tf, runs, size, color, *, font=None, bold=False,
              align=PP_ALIGN.LEFT, italic=False, space_after=6,
              first=False, bullet_char=None, bullet_color=None,
              indent_level=0):
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs \
            else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = 1.12
        if isinstance(runs, str):
            runs = [(runs, False)]
        if bullet_char:
            r = p.add_run()
            r.text = ("    " if indent_level else "") + bullet_char + "  "
            r.font.size = Pt(max(size - 4, 9))
            r.font.color.rgb = _rgb(bullet_color or self.c["accent"])
            r.font.name = self.f["body"]
        for text, run_bold in runs:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold or run_bold
            r.font.italic = italic
            r.font.color.rgb = _rgb(color)
            r.font.name = font or self.f["body"]
        return p

    def _card(self, slide, rect: Rect, fill=None, line=None):
        radius = self.rules.get("card_corner_radius", 0.1)
        shape_type = (MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0
                      else MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(
            shape_type, Inches(rect.left), Inches(rect.top),
            Inches(rect.width), Inches(rect.height))
        if radius > 0:
            try:
                shape.adjustments[0] = radius
            except (IndexError, ValueError):
                pass
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill or self.c["card_bg"])
        shape.line.color.rgb = _rgb(line or self.c["card_border"])
        shape.line.width = Pt(1)
        shape.shadow.inherit = False
        return shape

    def _footer(self, slide, page: int, deck_title: str):
        tf = self._textbox(
            slide, Rect(le.MARGIN, le.SLIDE_H - 0.42,
                        le.SLIDE_W - 2 * le.MARGIN, 0.3))
        p = self._para(tf, deck_title, SIZE["footer"], self.c["muted"],
                       first=True, space_after=0)
        r = p.add_run()
        r.text = f"   ·   {page:02d}"
        r.font.size = Pt(SIZE["footer"])
        r.font.color.rgb = _rgb(self.c["muted"])
        r.font.name = self.f["body"]

    def _slide_title(self, slide, text: str):
        tf = self._textbox(slide, le.title_rect(), anchor=MSO_ANCHOR.MIDDLE)
        size = autofit.fit_font_size(
            text, le.title_rect().width, le.title_rect().height,
            SIZE["title"], min_pt=20)
        self._para(tf, text, size, self.c["primary"],
                   font=self.f["title"], bold=True, first=True,
                   space_after=0)

    # -- hero (title / section / closing) ----------------------------------

    def render_hero(self, slide_model: Slide, kind: str, page: int):
        dark = self.rules.get("dark_title_slide", True)
        slide = self._new_slide(dark=dark)
        rects = le.hero_rects()
        title_color = self.c["text_inverse"] if dark else self.c["primary"]
        sub_color = self.c["text_inverse"] if dark else self.c["text"]

        # small kicker dot motif (theme accent), centered above the title
        dot_d = 0.16
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(le.SLIDE_W / 2 - dot_d / 2),
            Inches(rects["kicker"].top + 0.18), Inches(dot_d), Inches(dot_d))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(self.c["accent"])
        dot.line.fill.background()
        dot.shadow.inherit = False

        tf = self._textbox(slide, rects["title"], anchor=MSO_ANCHOR.MIDDLE)
        size = autofit.fit_font_size(
            slide_model.title, rects["title"].width, rects["title"].height,
            SIZE["hero_title"], min_pt=24)
        self._para(tf, slide_model.title, size, title_color,
                   font=self.f["title"], bold=True,
                   align=PP_ALIGN.CENTER, first=True, space_after=0)

        if slide_model.subtitle:
            tf = self._textbox(slide, rects["subtitle"])
            self._para(tf, slide_model.subtitle, SIZE["hero_sub"], sub_color,
                       align=PP_ALIGN.CENTER, italic=True, first=True,
                       space_after=0)

    # -- standard content ---------------------------------------------------

    def render_content(self, slide_model: Slide, page: int, deck: str):
        slide = self._new_slide()
        self._slide_title(slide, slide_model.title)
        area = le.body_rect()
        tf = self._textbox(slide, area)

        paras = []
        for b in slide_model.blocks:
            base = SIZE["col_header"] if b.kind == "heading" else SIZE["body"]
            paras.append((b.text, base))
        scale = autofit.fit_scale(paras, area.width, area.height)

        first = True
        for b in slide_model.blocks:
            if b.kind == "heading":
                self._para(tf, b.runs, max(int(SIZE["col_header"] * scale), 12),
                           self.c["primary"], font=self.f["title"], bold=True,
                           first=first, space_after=8)
            elif b.kind == "bullet":
                self._para(tf, b.runs, max(int(SIZE["body"] * scale), 10),
                           self.c["text"], first=first,
                           bullet_char=self.rules["bullet_char"],
                           indent_level=b.level)
            elif b.kind == "quote":
                self._para(tf, b.runs, max(int(SIZE["body"] * scale), 10),
                           self.c["secondary"], italic=True, first=first,
                           space_after=10)
            else:
                self._para(tf, b.runs, max(int(SIZE["body"] * scale), 10),
                           self.c["text"], first=first)
            first = False
        self._footer(slide, page, deck)

    # -- two column ----------------------------------------------------------

    def render_two_column(self, slide_model: Slide, page: int, deck: str):
        slide = self._new_slide()
        self._slide_title(slide, slide_model.title)

        # split blocks into the two ## sections; blocks appearing before
        # the first heading are folded into the first section so nothing
        # is silently dropped
        sections, current, orphans = [], None, []
        for b in slide_model.blocks:
            if b.kind == "heading":
                current = {"head": b, "items": []}
                sections.append(current)
            elif current is not None:
                current["items"].append(b)
            else:
                orphans.append(b)
        sections = sections[:2]
        if orphans and sections:
            sections[0]["items"][:0] = orphans

        for (head_rect, body_rect), sec in zip(le.column_rects(2), sections):
            # column header: bold colored text, no decorative bar
            tf = self._textbox(slide, head_rect, anchor=MSO_ANCHOR.MIDDLE)
            hsize = autofit.fit_font_size(
                sec["head"].text, head_rect.width, head_rect.height,
                SIZE["col_header"], min_pt=14)
            self._para(tf, sec["head"].runs, hsize, self.c["secondary"],
                       font=self.f["title"], bold=True, first=True,
                       space_after=0)

            card = self._card(slide, body_rect)
            tf = card.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = tf.margin_right = Inches(0.22)
            tf.margin_top = tf.margin_bottom = Inches(0.18)

            paras = [(b.text, SIZE["body"]) for b in sec["items"]]
            scale = autofit.fit_scale(
                paras, body_rect.width - 0.44, body_rect.height - 0.36)
            first = True
            for b in sec["items"]:
                self._para(tf, b.runs, max(int(SIZE["body"] * scale), 10),
                           self.c["text"], first=first,
                           bullet_char=(self.rules["bullet_char"]
                                        if b.kind == "bullet" else None),
                           indent_level=b.level)
                first = False
        self._footer(slide, page, deck)

    # -- cards (3-4 parallel bullets) -----------------------------------------

    @staticmethod
    def _split_card(block: Block) -> tuple[str, list[str]]:
        """bullet -> (card title, body lines)"""
        if block.children:
            return block.text, [c.text for c in block.children]
        for sep in ("：", ": ", " — ", " - "):
            if sep in block.text:
                head, _, rest = block.text.partition(sep)
                return head.strip(), [rest.strip()]
        return block.text, []

    def render_cards(self, slide_model: Slide, page: int, deck: str):
        slide = self._new_slide()
        self._slide_title(slide, slide_model.title)
        bullets = slide_model.top_bullets
        # quotes / paragraphs around the bullets render as a band below
        # the cards so no content is ever silently dropped
        extras = [b for b in slide_model.blocks
                  if b.kind in ("quote", "para")]
        reserve = 0.8 if extras else 0.0
        rects = le.card_rects(len(bullets), reserve_bottom=reserve)

        for i, (rect, block) in enumerate(zip(rects, bullets)):
            card = self._card(slide, rect)
            # numbered accent circle motif, top-left inside the card
            d = 0.42
            circ = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(rect.left + 0.22),
                Inches(rect.top + 0.22), Inches(d), Inches(d))
            circ.fill.solid()
            circ.fill.fore_color.rgb = _rgb(self.c["accent"])
            circ.line.fill.background()
            circ.shadow.inherit = False
            ctf = circ.text_frame
            ctf.margin_left = ctf.margin_right = 0
            ctf.margin_top = ctf.margin_bottom = 0
            p = ctf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(i + 1)
            r.font.size = Pt(15)
            r.font.bold = True
            r.font.color.rgb = _rgb(self.c["bg"])
            r.font.name = self.f["title"]

            title, lines = self._split_card(block)
            text_rect = Rect(rect.left + 0.24, rect.top + 0.22 + d + 0.12,
                             rect.width - 0.48,
                             rect.height - d - 0.7)
            tf = self._textbox(slide, text_rect)
            paras = [(title, SIZE["card_title"])] + \
                    [(ln, SIZE["card_body"]) for ln in lines]
            scale = autofit.fit_scale(paras, text_rect.width,
                                      text_rect.height)
            self._para(tf, title, max(int(SIZE["card_title"] * scale), 11),
                       self.c["primary"], font=self.f["title"], bold=True,
                       first=True, space_after=8)
            for ln in lines:
                self._para(tf, ln, max(int(SIZE["card_body"] * scale), 9),
                           self.c["text"], space_after=5)

        if extras:
            area = le.body_rect()
            band = Rect(area.left, area.top + area.height - reserve + 0.15,
                        area.width, reserve - 0.15)
            tf = self._textbox(slide, band, anchor=MSO_ANCHOR.MIDDLE)
            first = True
            for b in extras:
                size = autofit.fit_font_size(
                    b.text, band.width, band.height / len(extras),
                    SIZE["body"], min_pt=11)
                self._para(tf, b.runs, size, self.c["secondary"],
                           italic=(b.kind == "quote"), bold=True,
                           align=PP_ALIGN.CENTER, first=first,
                           space_after=2)
                first = False
        self._footer(slide, page, deck)

    # -- timeline / steps ------------------------------------------------------

    def render_timeline(self, slide_model: Slide, page: int, deck: str):
        slide = self._new_slide()
        self._slide_title(slide, slide_model.title)
        steps = slide_model.steps
        geo = le.timeline_positions(len(steps))

        # connector line first (sits behind the circles)
        for conn in geo["connectors"]:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(conn.left), Inches(conn.top),
                Inches(conn.width), Inches(conn.height))
            bar.fill.solid()
            bar.fill.fore_color.rgb = _rgb(self.c["card_border"])
            bar.line.fill.background()
            bar.shadow.inherit = False

        for i, (circle, label, step) in enumerate(
                zip(geo["circles"], geo["labels"], steps)):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(circle.left), Inches(circle.top),
                Inches(circle.width), Inches(circle.height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(
                self.c["primary"] if i % 2 == 0 else self.c["secondary"])
            shape.line.fill.background()
            shape.shadow.inherit = False
            ctf = shape.text_frame
            ctf.margin_left = ctf.margin_right = 0
            p = ctf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(i + 1)
            r.font.size = Pt(22)
            r.font.bold = True
            r.font.color.rgb = _rgb(self.c["bg"])
            r.font.name = self.f["title"]

            title, lines = self._split_card(step)
            tf = self._textbox(slide, label)
            paras = [(title, SIZE["step_label"])] + \
                    [(ln, SIZE["step_label"] - 2) for ln in lines]
            scale = autofit.fit_scale(paras, label.width, label.height)
            self._para(tf, title, max(int(SIZE["step_label"] * scale), 9),
                       self.c["primary"], bold=True, align=PP_ALIGN.CENTER,
                       first=True, space_after=4)
            for ln in lines:
                self._para(tf, ln,
                           max(int((SIZE["step_label"] - 2) * scale), 9),
                           self.c["muted"], align=PP_ALIGN.CENTER,
                           space_after=2)

        # non-step blocks (intro text, quotes) go into a bottom band
        extras = [b for b in slide_model.blocks if b.kind != "step"]
        if extras:
            band = Rect(le.MARGIN, le.SLIDE_H - le.MARGIN - 0.75,
                        le.SLIDE_W - 2 * le.MARGIN, 0.7)
            tf = self._textbox(slide, band, anchor=MSO_ANCHOR.MIDDLE)
            first = True
            for b in extras:
                size = autofit.fit_font_size(
                    b.text, band.width, band.height / len(extras),
                    SIZE["body"] - 2, min_pt=10)
                self._para(tf, b.runs, size, self.c["secondary"],
                           italic=(b.kind == "quote"),
                           align=PP_ALIGN.CENTER, first=first,
                           space_after=2)
                first = False
        self._footer(slide, page, deck)

    # -- entry point -------------------------------------------------------------

    def render(self, slides: list[Slide], out_path: str) -> str:
        deck_title = slides[0].title if slides else "SlideForge"
        total = len(slides)
        for i, sm in enumerate(slides):
            layout = le.detect_layout(sm, i, total)
            if layout in ("title", "closing", "section"):
                self.render_hero(sm, layout, i + 1)
            elif layout == "two_column":
                self.render_two_column(sm, i + 1, deck_title)
            elif layout == "cards":
                self.render_cards(sm, i + 1, deck_title)
            elif layout == "timeline":
                self.render_timeline(sm, i + 1, deck_title)
            else:
                self.render_content(sm, i + 1, deck_title)
        self.prs.save(out_path)
        return out_path
