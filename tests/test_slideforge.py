"""SlideForge test suite: parser, layout detection, autofit, end-to-end."""

import pathlib
import sys

import pytest

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parents[1]))

from slideforge import autofit, layout_engine as le  # noqa: E402
from slideforge.parser import parse_markdown, parse_runs  # noqa: E402
from slideforge.renderer import Renderer  # noqa: E402
from slideforge.sample import SAMPLE_MARKDOWN  # noqa: E402
from slideforge.themes import get_theme, list_themes  # noqa: E402


# -- parser -----------------------------------------------------------------

def test_h1_and_hr_split_slides():
    slides = parse_markdown("# A\n\ntext\n\n---\n\n# B\n- x\n")
    assert [s.title for s in slides] == ["A", "B"]
    assert slides[0].subtitle == "text"


def test_bold_runs():
    runs = parse_runs("normal **bold** tail")
    assert runs == [("normal ", False), ("bold", True), (" tail", False)]


def test_numbered_and_nested_bullets():
    slides = parse_markdown(
        "# S\n1. one\n2. two\n---\n# T\n- top\n  - child\n")
    assert len(slides[0].steps) == 2
    top = slides[1].top_bullets
    assert len(top) == 1 and top[0].children[0].text == "child"


# -- layout detection ---------------------------------------------------------

def _detect(md, index=1, total=5):
    slide = parse_markdown(md)[0]
    return le.detect_layout(slide, index, total)


def test_two_headings_give_two_column():
    assert _detect("# S\n## L\n- a\n## R\n- b\n") == "two_column"


def test_three_bullets_give_cards():
    assert _detect("# S\n- a: 1\n- b: 2\n- c: 3\n") == "cards"


def test_numbered_list_gives_timeline():
    assert _detect("# S\n1. a\n2. b\n3. c\n") == "timeline"


def test_first_and_last_title_only_are_hero():
    slide = parse_markdown("# Hello\nsub\n")[0]
    assert le.detect_layout(slide, 0, 5) == "title"
    assert le.detect_layout(slide, 4, 5) == "closing"


def test_default_is_content():
    assert _detect("# S\n- a\n- b\n") == "content"


# -- geometry -----------------------------------------------------------------

def test_card_rects_fit_canvas():
    for n in (3, 4):
        for r in le.card_rects(n):
            assert r.left >= le.MARGIN - 1e-6
            assert r.left + r.width <= le.SLIDE_W - le.MARGIN + 1e-6
            assert r.top + r.height <= le.SLIDE_H - le.MARGIN + 1e-6


def test_timeline_positions_counts():
    geo = le.timeline_positions(5)
    assert len(geo["circles"]) == len(geo["labels"]) == 5
    assert len(geo["connectors"]) == 1
    geo = le.timeline_positions(8)   # falls back to 2 rows
    assert len(geo["connectors"]) == 2


# -- autofit -------------------------------------------------------------------

def test_fit_font_size_shrinks_long_text():
    short = autofit.fit_font_size("Hi", 5.0, 1.0, 32)
    long = autofit.fit_font_size("word " * 120, 5.0, 1.0, 32)
    assert short == 32
    assert long < 32
    assert long >= autofit.MIN_SIZE
    # shrinks in 2pt steps unless it bottomed out at the minimum size
    assert long == autofit.MIN_SIZE or (32 - long) % autofit.STEP == 0


def test_fit_font_size_never_raises_on_extreme_text():
    size = autofit.fit_font_size("あ" * 5000, 1.0, 0.5, 40)
    assert size >= autofit.MIN_SIZE


def test_cjk_counts_wider_than_latin():
    assert autofit.text_width_em("あいう") > autofit.text_width_em("abc")


# -- themes ----------------------------------------------------------------------

def test_all_themes_have_required_keys():
    for name in list_themes():
        theme = get_theme(name)
        for key in ("primary", "secondary", "accent", "bg", "bg_dark",
                    "text", "text_inverse", "muted", "card_bg",
                    "card_border"):
            rgb = theme["colors"][key]
            assert len(rgb) == 3 and all(0 <= v <= 255 for v in rgb)
        assert theme["fonts"]["title"] and theme["fonts"]["body"]


def test_unknown_theme_raises():
    with pytest.raises(KeyError):
        get_theme("NoSuchTheme")


# -- end to end --------------------------------------------------------------------

@pytest.mark.parametrize("theme_name", list_themes())
def test_sample_renders_for_every_theme(tmp_path, theme_name):
    slides = parse_markdown(SAMPLE_MARKDOWN)
    out = tmp_path / f"sample_{theme_name}.pptx"
    Renderer(get_theme(theme_name)).render(slides, str(out))
    assert out.stat().st_size > 10_000

    # reopen and verify slide count + that text survived the round trip
    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(prs.slides) == len(slides)
    all_text = "\n".join(
        sh.text_frame.text for s in prs.slides
        for sh in s.shapes if sh.has_text_frame)
    assert "現状の課題" in all_text
    assert "AIレポート生成" in all_text
    # quotes attached to a cards slide must not be dropped
    assert "投資回収期間" in all_text


def test_no_content_lost_in_any_layout(tmp_path):
    """Every text line of the source markdown must appear in the deck."""
    from pptx import Presentation
    out = tmp_path / "full.pptx"
    slides = parse_markdown(SAMPLE_MARKDOWN)
    Renderer(get_theme("MinimalGray")).render(slides, str(out))
    prs = Presentation(str(out))
    def norm(s):
        # card splitting may drop the "：" separator between title/body,
        # so compare with separators and whitespace removed
        for ch in " 　：:—–-\n":
            s = s.replace(ch, "")
        return s

    deck_text = norm("".join(
        sh.text_frame.text for s in prs.slides
        for sh in s.shapes if sh.has_text_frame))
    for slide in slides:
        for block in slide.blocks:
            assert norm(block.text) in deck_text, f"missing: {block.text!r}"
